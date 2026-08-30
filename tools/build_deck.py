#!/usr/bin/env python3
"""Generate PRESENTATION.pptx for the Brand AI-Readiness Audit project.

Pure python-pptx + a matplotlib before/after chart. Run: python3 tools/build_deck.py
Output: PRESENTATION.pptx at the repo root.
"""
import os
import tempfile

os.environ.setdefault("MPLCONFIGDIR", tempfile.mkdtemp())

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
OUT = os.path.join(ROOT, "PRESENTATION.pptx")

NAVY = RGBColor(0x0B, 0x2E, 0x59)
BLUE = RGBColor(0x1F, 0x5C, 0x9E)
TEAL = RGBColor(0x18, 0x8A, 0x8A)
AMBER = RGBColor(0xC8, 0x7A, 0x00)
INK = RGBColor(0x22, 0x2A, 0x33)
GREY = RGBColor(0x5A, 0x63, 0x6E)
LIGHT = RGBColor(0xF3, 0xF5, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _box(slide, l, t, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill if fill else WHITE
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(2)
    for i, (txt, size, color, bold, indent) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.level = indent
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def bullet_slide(title, bullets, kicker="", accent=BLUE):
    """bullets: list of (text, level) — level 0 main, 1 sub."""
    s = prs.slides.add_slide(BLANK)
    _box(s, 0, 0, SW, Inches(1.15), fill=NAVY)
    _box(s, 0, Inches(1.15), SW, Pt(4), fill=accent)
    if kicker:
        _text(s, Inches(0.6), Inches(0.16), Inches(12), Inches(0.3),
              [(kicker.upper(), 12, TEAL, True, 0)])
    _text(s, Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.7),
          [(title, 27, WHITE, True, 0)], anchor=MSO_ANCHOR.MIDDLE)
    runs = []
    for txt, lvl in bullets:
        if lvl == 0:
            runs.append(("• " + txt, 18, INK, False, 0))
        else:
            runs.append(("– " + txt, 15, GREY, False, 1))
    _text(s, Inches(0.75), Inches(1.5), Inches(11.9), Inches(5.6), runs, space=9)
    return s


def table_slide(title, headers, rows, kicker="", col_widths=None, accent=BLUE, fontsize=13):
    s = prs.slides.add_slide(BLANK)
    _box(s, 0, 0, SW, Inches(1.15), fill=NAVY)
    _box(s, 0, Inches(1.15), SW, Pt(4), fill=accent)
    if kicker:
        _text(s, Inches(0.6), Inches(0.16), Inches(12), Inches(0.3),
              [(kicker.upper(), 12, TEAL, True, 0)])
    _text(s, Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.7),
          [(title, 27, WHITE, True, 0)], anchor=MSO_ANCHOR.MIDDLE)
    nrows, ncols = len(rows) + 1, len(headers)
    tw = Inches(12.1); th = Inches(0.5) * nrows
    gt = s.shapes.add_table(nrows, ncols, Inches(0.6), Inches(1.5), tw, th).table
    if col_widths:
        for i, w in enumerate(col_widths):
            gt.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        c = gt.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = BLUE
        c.margin_left = c.margin_right = Pt(6)
        p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = h
        r.font.bold = True; r.font.size = Pt(fontsize + 1); r.font.color.rgb = WHITE; r.font.name = "Calibri"
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.fill.solid()
            c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
            c.margin_left = c.margin_right = Pt(6); c.margin_top = c.margin_bottom = Pt(3)
            p = c.text_frame.paragraphs[0]; p.word_wrap = True
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(fontsize); r.font.color.rgb = INK; r.font.name = "Calibri"
            if j == 0:
                r.font.bold = True; r.font.color.rgb = NAVY
    return s


def make_chart():
    path = os.path.join(tempfile.gettempdir(), "beforeafter.png")
    fig, (a1, a2) = plt.subplots(1, 2, figsize=(9.6, 4.2))
    fig.patch.set_facecolor("white")
    labels = ["Round 1\n(pre-fix)", "Held-out\n(post-fix)"]
    navy = "#1F5C9E"; teal = "#188A8A"; amber = "#C87A00"

    a1.bar(labels, [0.576, 0.857], color=[amber, teal], width=0.55)
    a1.set_ylim(0, 1.0); a1.set_title("Precision (blind eval)", fontsize=13, weight="bold", color="#0B2E59")
    for i, v in enumerate([0.576, 0.857]):
        a1.text(i, v + 0.02, f"{v:.3f}", ha="center", fontsize=13, weight="bold", color="#222")
    a1.spines[["top", "right"]].set_visible(False)
    a1.tick_params(labelsize=11)

    a2.bar(labels, [36, 5], color=[amber, teal], width=0.55)
    a2.set_ylim(0, 40); a2.set_title("False positives", fontsize=13, weight="bold", color="#0B2E59")
    for i, v in enumerate([36, 5]):
        a2.text(i, v + 0.8, str(v), ha="center", fontsize=13, weight="bold", color="#222")
    a2.spines[["top", "right"]].set_visible(False)
    a2.tick_params(labelsize=11)

    fig.tight_layout(pad=1.5)
    fig.savefig(path, dpi=160, bbox_inches="tight")
    plt.close(fig)
    return path


def chart_slide(title, img, note, kicker=""):
    s = prs.slides.add_slide(BLANK)
    _box(s, 0, 0, SW, Inches(1.15), fill=NAVY)
    _box(s, 0, Inches(1.15), SW, Pt(4), fill=TEAL)
    if kicker:
        _text(s, Inches(0.6), Inches(0.16), Inches(12), Inches(0.3),
              [(kicker.upper(), 12, TEAL, True, 0)])
    _text(s, Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.7),
          [(title, 27, WHITE, True, 0)], anchor=MSO_ANCHOR.MIDDLE)
    s.shapes.add_picture(img, Inches(1.7), Inches(1.55), width=Inches(10))
    _text(s, Inches(0.75), Inches(6.55), Inches(11.9), Inches(0.7),
          [(note, 16, TEAL, True, 0)], align=PP_ALIGN.CENTER)
    return s


def title_slide():
    s = prs.slides.add_slide(BLANK)
    _box(s, 0, 0, SW, SH, fill=NAVY)
    _box(s, 0, Inches(3.05), SW, Pt(4), fill=TEAL)
    _text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.0),
          [("Brand AI-Readiness Audit", 44, WHITE, True, 0)])
    _text(s, Inches(0.9), Inches(3.25), Inches(11.5), Inches(1.2),
          [("An Agent Skill Marketplace that audits any website for AI discoverability "
            "and on-site engagement — and proves it works on sites it has never seen.",
            20, RGBColor(0xC9, 0xD6, 0xE6), False, 0)])
    _text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
          [("Recommend-only · read-only · deterministic · < 5 min · generalizes by construction",
            14, TEAL, True, 0)])
    return s


def closing_slide():
    s = prs.slides.add_slide(BLANK)
    _box(s, 0, 0, SW, SH, fill=NAVY)
    _text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.0),
          [("Precision 0.58 → 0.86 on unseen sites,", 32, WHITE, True, 0),
           ("measured by a test harness that never grades itself.", 32, TEAL, True, 0)])
    _text(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(1.4),
          [("6 agentskills.io-valid skills · 31 checks · one schema-valid report", 16, RGBColor(0xC9, 0xD6, 0xE6), False, 0),
           ("Blind red-team dataset + independent adjudication + debate → generalizable fixes", 16, RGBColor(0xC9, 0xD6, 0xE6), False, 0)])
    return s


# ---- Build the deck ----
title_slide()

bullet_slide("The challenge", [
    ("Point an agent at ANY website; audit it and emit one findings-plus-fixes report.", 0),
    ("Cover both halves of the problem:", 0),
    ("AI discoverability — why AI assistants don't find or cite the brand", 1),
    ("On-site engagement — why visitors who arrive don't stay", 1),
    ("Recommend-only, read-only, deterministic, < 5 min, no site changes.", 0),
    ("Graded on the skill's LOGIC and on GENERALIZATION — no example sites are given.", 0),
], kicker="Round 3 — Agent Skill Marketplace")

bullet_slide("Why it's really two problems", [
    ("An AI assistant only cites a page it can:", 0),
    ("reach (crawler let in) → read (no JS wall) → pick a clean fact → trust it (entity + corroboration)", 1),
    ("A visitor only stays on a page that loads fast and orients them quickly.", 0),
    ("Two distinct failure surfaces → the marketplace is split along the real mechanism chain.", 0),
    ("Every finding is tagged discoverability / engagement / both.", 0),
], kicker="The mechanism")

bullet_slide("What we built", [
    ("One entrypoint skill orchestrating five focused sub-skills over a single shared, read-only crawl.", 0),
    ("Hybrid: deterministic scripts do the mechanical signals; the AI agent does the judgment + fixes.", 0),
    ("Shared common/ spine: fetch + optional render, robots, finding contract, report assembly.", 0),
    ("31 checks across the two halves, each with explicit false-positive guards.", 0),
    ("Each sub-skill is a valid agentskills.io skill (SKILL.md + scripts/ + references/) and runs on its own.", 0),
    ("marketplace.json names exactly one entrypoint; self-contained (1.8 MB, no model weights).", 0),
], kicker="Architecture")

table_slide("The five sub-skills (+ one entrypoint)",
            ["Skill", "The question it answers"],
            [["website-ai-audit", "Entrypoint: shared crawl → runs the 5 sub-skills → merges one schema-valid report"],
             ["crawler-access", "Can an AI citation crawler get in? robots per-bot, WAF/UA, indexability, sitemap"],
             ["machine-readability", "Once in, can a non-rendering crawler read it? JS render gap, latency"],
             ["fact-extractability", "Can a machine pick out the fact? JSON-LD, semantic HTML, facts-in-images"],
             ["entity-trust", "Is the fact trusted + attributed right? Organization anchor, sameAs, freshness"],
             ["onsite-engagement", "Why don't visitors stay? speed, a11y, interstitials, mixed content, readability"]],
            kicker="Separation of concerns", col_widths=[3.0, 9.1], fontsize=12.5)

bullet_slide("Hybrid: deterministic signals + agent judgment", [
    ("Agent Skills are run BY an AI agent — so we split the work by what each side is best at.", 0),
    ("Scripts own the mechanical, measurable signals (robots, render gap, JSON-LD, headers, contrast).", 0),
    ("Reproducible, no LLM — this is the layer whose precision we measure.", 1),
    ("The agent owns judgment: answer-readiness, entity clarity, value-prop — and site-specific fixes.", 0),
    ("It reasons over an evidence pack the scripts emit; compose.py merges both into one report.", 1),
    ("Runs headless too (scripts alone = full report, no LLM/gateway); the agent layer is additive.", 0),
], kicker="Execution model", accent=TEAL)

bullet_slide("One polite, read-only pass", [
    ("Deterministic scope: homepage + a sorted, capped sample of internal/sitemap URLs.", 0),
    ("One shared fetch (and optional render) per page; sub-skills read the cache — never re-fetch.", 0),
    ("robots.txt honored; GET/HEAD only; no auth; nothing written to the site.", 0),
    ("Playwright render when available; otherwise a static SPA heuristic — it never crashes.", 0),
    ("Byte-identical output across runs (timestamp injectable); ~0.03 s per fixture audit.", 0),
], kicker="How it runs")

bullet_slide("What a finding looks like", [
    ("id · title · severity (critical/high/medium/low)", 0),
    ("evidence — concrete, with numbers + prevalence: \"0/12 product pages carry schema.org markup\"", 0),
    ("category · half · the mechanism it addresses", 0),
    ("suggested_action — fix summary, a fix-first priority (distinct from severity), effort, confidence", 0),
    ("Report adds counts-by-severity, a discoverability/engagement split, and an honest limitations block.", 0),
    ("Validated against a bundled JSON Schema before emit; counts computed, never hand-kept.", 0),
], kicker="Output design")

bullet_slide("Built-in guarantees", [
    ("Read-only & recommend-only — never modifies the audited site.", 0),
    ("Deterministic — fixed sampling, sorted output, bucketed thresholds.", 0),
    ("Fast & bounded — one shared pass, page cap, per-check timeouts, well under 5 min.", 0),
    ("Portable & small — pure-Python core, 1.8 MB, no model weights; browser/OCR optional.", 0),
    ("Conservative — every check has FP guards; blocked/unreadable pages never yield 'content missing'.", 0),
], kicker="Safety & hygiene")

bullet_slide("Testing honestly: the bias problem", [
    ("If the same author writes the detector, the tests, AND the grader, the tests prove nothing.", 0),
    ("Our first fixtures — whose authors had seen the checks — reported a flawless 0% false-positive rate.", 0),
    ("That number was meaningless. So we rebuilt testing around a strict information barrier.", 0),
    ("Principle: no single model both creates the ground truth and judges against it.", 0),
], kicker="The journey — why", accent=AMBER)

table_slide("The information barrier",
            ["Role", "Who", "Sees the detector?", "Model"],
            [["System under test", "the auditor (deterministic Python)", "it IS the detector", "code"],
             ["Test-set author", "blind red-team generators", "no — walled off", "Sonnet"],
             ["Grader", "blind adjudicators", "no — judge on merits", "Opus"],
             ["Tie-breakers", "prosecutor / defender / referee", "no", "Sonnet / Opus"]],
            kicker="The journey — design", col_widths=[2.6, 4.6, 3.0, 1.9], accent=AMBER, fontsize=12.5)

bullet_slide("Step 1 — a ground-truth dataset, kept blind", [
    ("An independent red team authored realistic sites: SaaS SPA, restaurant, e-commerce, news,", 0),
    ("corporate brand, personal blog, government page, and deliberately-clean controls.", 1),
    ("For each site it sealed its OWN ground truth: every real problem (severity + mechanism),", 0),
    ("plus what's FINE and the legitimate TRAPS that must NOT be flagged.", 1),
    ("The authors never saw how the auditor detects anything → trustworthy, unbiased labels.", 0),
], kicker="The journey — dataset", accent=AMBER)

bullet_slide("Step 2 — the first honest measurement", [
    ("Ran the deterministic auditor over the blind sites.", 0),
    ("Independent judges inspected each site themselves and scored every finding TP vs FP.", 0),
    ("A prosecutor/defender/referee debate settled every contested finding.", 0),
    ("Verdict: precision 0.576 — 49 true positives against 36 false positives.", 0),
    ("The biased fixtures had hidden all 36.", 0),
], kicker="The journey — measure", accent=AMBER)

bullet_slide("Step 3 — what the judges + debate found", [
    ("Flagging 'plain HTTP' when the site's own canonical/sitemap declare HTTPS", 0),
    ("Over-claiming 'render-blocking/heavy' on tiny pages (only real factor: no gzip)", 0),
    ("Penalizing a legitimate off-domain canonical (staging/CDN)", 0),
    ("Treating ClaudeBot / Meta-ExternalAgent (training crawlers) as citation bots", 0),
    ("Rejecting a valid ProfilePage/Person author page; missing nested publisher Organization", 0),
    ("Auditing robots.txt as a page; emitting noise for a missing llms.txt", 0),
], kicker="The journey — root causes", accent=AMBER)

bullet_slide("Step 4 — fix generalizably", [
    ("Two fix cycles corrected the LOGIC, not the symptoms.", 0),
    ("Each fixer was told the general bug + mechanism — and forbidden from special-casing any site.", 0),
    ("Result on the round-1 set: 28 false positives removed, zero real regressions.", 0),
    ("White-box regression stayed perfect (26/26, 0 FP); harnesses green.", 0),
], kicker="The journey — fix", accent=AMBER)

chart_slide("Step 5 — proof it generalizes", make_chart(),
            "Held-out sites were generated AFTER and blind to every fix → generalization, not overfitting.",
            kicker="The journey — held-out")

table_slide("Results against the rubric",
            ["Rubric criterion", "How we meet it"],
            [["Detection accuracy / few FPs / both halves", "0.86 precision on UNSEEN sites (blind eval), per-half coverage"],
             ["Suggested-action quality", "mechanism + fix-first priority + effort + confidence per finding"],
             ["Output design", "one schema-validated JSON report + readable Markdown view"],
             ["Skill / marketplace hygiene", "6 valid skills, exactly 1 entrypoint, deterministic, safe"],
             ["Generalization", "headline number comes from held-out sites made after the fixes"]],
            kicker="Scorecard", col_widths=[4.6, 7.5], fontsize=12.5)

bullet_slide("Limitations & next steps", [
    ("Recall is the open axis — 28 lower-severity misses on the held-out set (meta descriptions,", 0),
    ("dark-bg contrast, dead CTAs, per-page social cards). Same blind loop applies next.", 1),
    ("Live network is out of scope in the sandbox → field suite runs via offline snapshots (reproducible).", 0),
    ("A few checks are network-dependent (sameAs, broken links) and are labeled as such.", 0),
    ("Next: recall pass, real-site snapshot capture, and a live demo run.", 0),
], kicker="Honest gaps")

closing_slide()

prs.save(OUT)
print(f"wrote {OUT} — {len(prs.slides._sldIdLst)} slides")
